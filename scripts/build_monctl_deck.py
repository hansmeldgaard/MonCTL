"""Generate a technical overview PowerPoint for MonCTL.

Produces `MonCTL-technical-overview.pptx` in the repo root. Content is
sourced from `CLAUDE.md` and the current state of the codebase — every
claim here matches a path / feature that actually exists.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

# ── Design tokens ──────────────────────────────────────────
BRAND = RGBColor(0x2E, 0x7D, 0x64)        # muted teal — MonCTL brand feel
DARK_BG = RGBColor(0x15, 0x1A, 0x1F)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x52, 0x60, 0x72)
ACCENT = RGBColor(0x29, 0xA3, 0x85)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _text_frame(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    return tf


def _add_text_box(
    slide, left, top, width, height, text, *, size=18, bold=False,
    color=TEXT_DARK, font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = _text_frame(box)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(
    slide, left, top, width, height, bullets: list[str], *,
    size=16, color=TEXT_DARK, font="Calibri",
):
    """Add a bulleted list. Each bullet is rendered as its own paragraph
    with a leading dot and 12pt line spacing for readability."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = _text_frame(box)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def _add_header(slide, title: str, subtitle: str | None = None):
    """Standard slide header: brand stripe + title + optional subtitle."""
    # Top brand stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BRAND
    stripe.line.fill.background()

    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8),
        title, size=32, bold=True, color=TEXT_DARK,
    )
    if subtitle:
        _add_text_box(
            slide, Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.45),
            subtitle, size=16, color=TEXT_MUTED,
        )

    # Footer rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55),
        Inches(1.2), Inches(0.04),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def _add_footer(slide, idx: int, total: int):
    _add_text_box(
        slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
        "MonCTL — technical overview", size=10, color=TEXT_MUTED,
    )
    _add_text_box(
        slide, Inches(11.8), Inches(7.05), Inches(1.5), Inches(0.3),
        f"{idx} / {total}", size=10, color=TEXT_MUTED,
    )


# ── Slide builders ─────────────────────────────────────────
def make_title_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BG)

    # Brand bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.2), SLIDE_W, Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_text_box(
        slide, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
        "MonCTL", size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    _add_text_box(
        slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.7),
        "A distributed monitoring platform",
        size=28, color=RGBColor(0xBB, 0xC8, 0xD2),
    )
    _add_text_box(
        slide, Inches(0.7), Inches(4.2), Inches(12), Inches(0.5),
        "Technical overview",
        size=18, color=RGBColor(0x8A, 0x9A, 0xA8),
    )
    _add_text_box(
        slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
        "Central HA · Distributed collectors · ClickHouse time-series",
        size=14, color=RGBColor(0x6A, 0x7A, 0x88),
    )


def make_problem_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Problem & Mission",
        "Observability for heterogeneous network + system estates",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Fleets span Cisco, Juniper, Linux, Docker hosts, HTTP services — no single-vendor tool fits.",
            "Engineers need per-device metrics history, availability, interface counters, and alerts in one pane.",
            "Scale is horizontal: central control plane must tolerate node failures; collectors must partition jobs automatically.",
            "Operators need extensibility — custom monitoring \"apps\" and \"connectors\" without a full release cycle.",
            "Alerts must promote into incidents with escalation, suppression, and automation-driven run books.",
            "Everything driven by a first-class REST API + a responsive React UI for operators.",
        ],
        size=18,
    )
    _add_footer(slide, idx, total)


def make_architecture_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(slide, "High-level architecture", "Edge collectors pull jobs, push results to a clustered central")

    # Simple boxes representing topology
    def box(x, y, w, h, label, fill=BRAND, sub=""):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = ACCENT
        tf = _text_frame(shp)
        p = tf.paragraphs[0]
        p.alignment = 2  # center
        run = p.add_run()
        run.text = label
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = 2
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(11)
            r2.font.color.rgb = RGBColor(0xE0, 0xE8, 0xED)

    # Users / browser
    box(Inches(0.8), Inches(2.2), Inches(2.2), Inches(0.9),
        "Browser / API", sub="JWT cookies · Bearer keys",
        fill=RGBColor(0x3D, 0x5A, 0x80))

    # HAProxy VIP
    box(Inches(3.4), Inches(2.2), Inches(2.6), Inches(0.9),
        "HAProxy + keepalived",
        sub="VIP 10.145.210.40 :443", fill=RGBColor(0x29, 0x54, 0x70))

    # Central cluster
    box(Inches(6.4), Inches(1.8), Inches(3.2), Inches(1.6),
        "Central cluster (x4)",
        sub="FastAPI · Scheduler · Lifecycle\ncentral1-4", fill=BRAND)

    # Data stores
    box(Inches(10.1), Inches(1.5), Inches(2.8), Inches(0.8),
        "PostgreSQL (Patroni)", fill=RGBColor(0x44, 0x6B, 0x8A))
    box(Inches(10.1), Inches(2.4), Inches(2.8), Inches(0.8),
        "ClickHouse (replicated)", fill=RGBColor(0x44, 0x6B, 0x8A))
    box(Inches(10.1), Inches(3.3), Inches(2.8), Inches(0.8),
        "Redis + Sentinel", fill=RGBColor(0x44, 0x6B, 0x8A))

    # Collectors
    box(Inches(2), Inches(5.0), Inches(3), Inches(1.2),
        "Collectors (worker1-4)",
        sub="BasePoller · gRPC peer mesh\nResult buffer (SQLite)",
        fill=RGBColor(0x5C, 0x3A, 0x6B))

    # Devices
    box(Inches(6.5), Inches(5.0), Inches(5.2), Inches(1.2),
        "Monitored devices",
        sub="SNMP · ICMP · TCP · HTTP · SSH · Docker",
        fill=RGBColor(0x3A, 0x4B, 0x5C))

    _add_text_box(
        slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
        "Collectors poll central for job assignments, run checks, and forward results; "
        "central persists to ClickHouse (timeseries) + PostgreSQL (config).",
        size=13, color=TEXT_MUTED,
    )
    _add_footer(slide, idx, total)


def make_stack_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(slide, "Tech stack", "Python/FastAPI backend · React/TS frontend · Postgres + ClickHouse + Redis")

    rows = [
        ("Backend", "Python 3.11 · FastAPI · SQLAlchemy async · Alembic"),
        ("Frontend", "React 19 · TypeScript · Vite · Tailwind v4 · React Query"),
        ("Charts / icons", "Recharts · Lucide React · react-grid-layout"),
        ("Time-series", "ClickHouse (replicated cluster) — 4 domain tables + alert_log + events"),
        ("Relational", "PostgreSQL 16 with Patroni for HA + etcd DCS"),
        ("Cache / broker", "Redis Sentinel (announce-ip configured)"),
        ("Proxy / LB", "HAProxy + keepalived (VIP 10.145.210.40)"),
        ("Collector runtime", "Python 3.12 · gRPC (cache-node ↔ worker) · SQLite buffer"),
        ("Deploy", "Docker Compose per role · ./deploy.sh parallel rollout"),
    ]
    top = Inches(1.9)
    for i, (label, value) in enumerate(rows):
        _add_text_box(
            slide, Inches(0.7), top + Inches(i * 0.5), Inches(3), Inches(0.45),
            label, size=16, bold=True, color=ACCENT,
        )
        _add_text_box(
            slide, Inches(3.8), top + Inches(i * 0.5), Inches(9), Inches(0.45),
            value, size=16, color=TEXT_DARK,
        )
    _add_footer(slide, idx, total)


def make_central_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Central server", "FastAPI monolith with ~30 routed domains",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(6), Inches(5),
        [
            "Mount point: /v1/* (UI) · /api/v1/* (collectors)",
            "Auth: JWT HTTP-only cookies + Bearer API keys",
            "Dependencies: get_db · get_clickhouse · require_permission",
            "Leader-elected scheduler: Redis SETNX 30s TTL",
            "Background: rebalancer · alert engine · incident engine",
            "ClickHouse batching buffer (monctl_central.cache)",
            "WebSocket command channel to collectors",
        ],
        size=15,
    )
    _add_bullets(
        slide, Inches(7), Inches(1.9), Inches(5.8), Inches(5),
        [
            "Domains: devices · apps · assignments · collectors",
            "credentials · connectors · templates · packs",
            "alerting · incidents · incident_rules · automations",
            "results · config_history · dashboard · analytics",
            "docker_infra · logs · audit · python_modules",
            "system · settings · tenants · users · roles",
            "ws · retention · upgrades · discovery",
        ],
        size=15, color=TEXT_MUTED,
    )
    _add_footer(slide, idx, total)


def make_collector_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Collectors",
        "Pull-based workers with local result buffer and gRPC peer mesh",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "api_client pulls job assignments from /api/v1/jobs on a delta-sync basis (updated_at cursor).",
            "Polling engine (min-heap deadline scheduler) schedules each assignment at its own interval.",
            "Connector lifecycle is engine-owned: connect → use → close. Apps MUST NOT call connect().",
            "Result forwarder batches PollResults and posts them back to central; SQLite buffer survives central outages.",
            "Persistent WebSocket command channel: poll_device, config_reload, docker_* — instantaneous ops from UI.",
            "Log shipper tails container logs → POST /v1/logs with configurable level filter per collector.",
            "gRPC peer (cache-node ↔ worker): shared app cache, reduces cold-start app downloads across nodes.",
            "Memory hygiene: gc.collect() + malloc_trim(0) after every job; MALLOC_ARENA_MAX=2, mem_limit: 1g.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_ha_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "High availability",
        "Every layer is clustered — no single point of failure",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "HAProxy + keepalived on all 4 central nodes · VIP 10.145.210.40 floats on failover.",
            "TLS terminated at HAProxy; round-robin to 4 central app instances on :8443.",
            "PostgreSQL HA via Patroni on central1+2 with etcd (3 nodes) as DCS.",
            "ClickHouse replicated cluster on central3+4 with a Keeper quorum on central1.",
            "Redis Sentinel (3 nodes) — announce-ip mandatory in Docker bridge mode; coordinated restart to fix state.",
            "Scheduler leader election: Redis SETNX with 30s TTL, renewed every 10s — only one instance runs scheduled tasks.",
            "Deploy model: ./deploy.sh pushes one saved image to all 4 nodes in parallel (~20s rollout).",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_storage_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Data stores & retention",
        "Right tool for each workload",
    )
    rows = [
        ("PostgreSQL", "All operator config — devices, apps, assignments, credentials, users, incidents, automations. SQLAlchemy async ORM."),
        ("ClickHouse", "4 domain result tables (availability_latency, performance, config, interface) + alert_log + events + logs."),
        ("Retention tiers", "Raw 7d → hourly 90d → daily 730d. Tier auto-selected by query time range."),
        ("TTL gotcha", "CH 24.3 can't parse ISO; timestamps wrapped as YYYY-MM-DD HH:MM:SS; toDateTime() wrapper in TTL for DateTime64."),
        ("Audit split", "audit_login_events in PG (ACID) · audit_mutations in ClickHouse (365d TTL, high-volume)."),
        ("Redis", "Rate counters, discovery flags, poll-now pubsub, sentinel-managed failover. Sessions and caches."),
        ("SQLite (collector)", "Local result buffer; survives central outages without data loss."),
    ]
    top = Inches(1.9)
    for i, (label, value) in enumerate(rows):
        _add_text_box(
            slide, Inches(0.7), top + Inches(i * 0.65), Inches(2.7), Inches(0.5),
            label, size=16, bold=True, color=ACCENT,
        )
        _add_text_box(
            slide, Inches(3.5), top + Inches(i * 0.65), Inches(9.5), Inches(0.9),
            value, size=14, color=TEXT_DARK,
        )
    _add_footer(slide, idx, total)


def make_scheduling_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Job scheduling & partitioning",
        "Consistent hashing with EMA-blended weight rebalancing",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Unpinned assignments partitioned via consistent hashing over CollectorGroup.weight_snapshot (JSONB).",
            "weight_snapshot = NULL → equal weights fallback. Reset on DOWN / approve / group-change.",
            "Rebalancer cycles every 5 minutes · imbalance threshold 1.5x · 15% hysteresis on weight change.",
            "Weight updates use EMA blending (α=0.3). Pure 1/cost weights flip the entire distribution; EMA converges over 5–7 cycles.",
            "Minimum weight floor = 0.3 — lower values produce extreme vnode ratios (20:1) and flip-flopping.",
            "Collectors identify themselves via MONCTL_COLLECTOR_ID — missing id = gets ALL jobs (guardrail, not a feature).",
            "Poll Now uses Redis pub/sub broadcast to all active collectors — instant, no round-trip via DB.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_apps_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Apps & packs",
        "BasePoller subclasses distributed as version-pinned wheels",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "All apps subclass BasePoller with async def poll() -> PollResult. No legacy stdin/stdout scripts.",
            "Built-in packs (packs/) auto-imported at startup using \"skip\" resolution — never overwrites existing.",
            "Current packs: snmp-core (snmp_check, snmp_interface_poller, snmp_discovery, snmp_uptime) · basic-checks (ping, port, http).",
            "snmp_discovery is a one-shot (interval=0) injected when Redis discovery flag set. Looked up by name.",
            "Auto-rename: devices bulk-imported with IP-as-name get renamed to sysName on discovery (hostname becomes primary identity).",
            "Error categories — every app PollResult MUST set error_category: device / config / app / \"\". Drives alerting.",
            "Versioning: app wheels distributed via central PyPI mirror · venv-per-app with permanent site-packages on sys.path.",
            "Connector slots (feat branch, pending): alias → slot model to decouple app code from credential shapes.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_credentials_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Connectors & credential resolution",
        "Three-level precedence gives per-assignment overrides without duplication",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Credential resolution chain (highest wins):",
            "   1. AssignmentCredentialOverride — per-assignment, per-connector-alias",
            "   2. AppAssignment.credential_id — assignment-level",
            "   3. Device.credentials JSONB — per-credential-type default for the device",
            "Connector lifecycle is engine-managed: engine calls .connect() once per job, .close() after. Idempotent guard in connectors.",
            "SNMP connector accepts both snmp_version and version keys — backward-compat fix (v1.3.1).",
            "Credential values encrypted at rest with app-managed key; field value never leaves central in plaintext.",
            "Reusable credential templates keyed by type (SSH, SNMP, HTTP Basic, etc.) drive UI + connector validation.",
            "Credential cache on collectors encrypted locally; rotated on WebSocket credential_invalidate command.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_interface_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Interface monitoring",
        "Targeted polling + central-side rate calculation",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Full walk (ifTable + ifXTable) on discovery · targeted GET thereafter using monitored_interfaces param.",
            "Rate calculation happens in central — previous counters cached in Redis, not on the collector.",
            "Multi-tier ClickHouse retention: raw 7d → hourly 90d → daily 730d. Tier auto-selected by query window.",
            "Interface metadata (description, admin/oper status, speed) tracked in PG for UI surfaces + filters.",
            "History default /v1/results queries availability_latency, performance, config — NOT interface (has its own endpoints).",
            "formatThresholdValue auto-scales bps → Mbps, bytes → GB for the threshold UI.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_alerting_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Alerting",
        "Compile-to-SQL DSL + 4-level threshold hierarchy",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "DSL examples: avg(metric) > 80 · rate(octets)*8/1e6 > 500 · field CHANGED · state IN (...)",
            "Expressions compile to ClickHouse SQL at evaluation time.",
            "Threshold hierarchy (highest wins): expression default → app value → device override → entity override.",
            "Engine runs on a 30s cycle; writes fire/clear transitions to the alert_log CH table.",
            "AlertEntity rows live in PG with lifecycle (firing, resolved, disabled) + fire_count + fire_history window.",
            "Labels drive downstream routing — device_id, device_name, assignment_id, and domain-specific keys.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_incidents_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Incident engine",
        "Replaced the legacy EventPolicy system across 14 PRs (#20–#33)",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "IncidentRule + Incident (PG) — operator-configured promotion of alerts → tracked incidents.",
            "Modes: consecutive or cumulative over a window_size. Thresholds gate when an incident OPENS.",
            "Severity ladder: wall-clock escalation from opened_at — never demotes. Validator enforces monotonic priority.",
            "Scope filter: AND-match against AlertEntity.entity_labels · empty string = wildcard (key must exist).",
            "depends_on + match_on: suppress child incidents under open parents. parent_incident_id is the suppression link.",
            "Flap guard: hold auto-clear for N seconds after alert resolves; re-fire keeps the same incident — no UUID churn.",
            "Companion-alert clear (latest): clear incident when ANY or ALL of a set of companion alerts are healthy for the same entity.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_automations_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Automations (Run Book Automation)",
        "Trigger actions on schedules or on incident transitions",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Two trigger types today: cron schedule · on_incident_transition (opened / escalated / cleared / any).",
            "Filters: incident_rule_ids (subset of rules) · incident_state_trigger (which transitions).",
            "Actions composed from reusable Action definitions — shell commands, SSH runbooks, HTTP callouts.",
            "Cooldown prevents runaway automations on flapping alerts; last-run time tracked in ClickHouse.",
            "AutomationEngine.on_incident_transition fires after IncidentEngine session.commit() — dispatched with plain dicts (no ORM refs).",
            "The legacy on_new_events hook was fully retired in cut-over step 2 (PR #31) — trigger_type='event' now rejected at the API.",
        ],
        size=16,
    )
    _add_footer(slide, idx, total)


def make_security_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Security & RBAC",
        "JWT cookies for the UI, Bearer keys for the API, custom roles for everything",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Web UI: access_token + refresh_token HTTP-only cookies; CSRF protection on state-changing routes.",
            "Management API: Authorization: Bearer monctl_m_<key> · hashed at rest (monctl_common.hash_api_key).",
            "Collector API: shared secret MONCTL_COLLECTOR_API_KEY on /api/v1/* endpoints.",
            "RBAC resources: device · app · assignment · credential · alert · collector · tenant · user · template · settings · result.",
            "Actions: view · create · edit · delete · manage. Admins and API keys bypass the permission check.",
            "Frontend: usePermissions() drives sidebar visibility, settings tabs, and action buttons — no UI reveal of unauthorized routes.",
            "Multi-tenant: apply_tenant_filter(stmt, auth, Device.tenant_id) scopes every list query.",
            "TLS terminated at HAProxy with keepalived-managed VIP — one certificate path, live sync via /settings/tls.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_audit_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Audit log",
        "Two-store design: PG for auth, ClickHouse for mutations",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "audit_login_events (PostgreSQL): every login success/failure with IP, user-agent, request_id. ACID.",
            "audit_mutations (ClickHouse): create/update/delete on whitelisted tables — 365d TTL, high-volume.",
            "Capture mechanism: AuditContextMiddleware puts user/IP/request_id in a contextvars.ContextVar.",
            "SQLAlchemy before_flush listener walks session.new/dirty/deleted, extracts old values from attrs.history.",
            "Whitelist (not blacklist): audit/resource_map.py::TABLE_TO_RESOURCE — adding a new audited table requires updating it.",
            "Redaction: fields matching password/api_key/token/secret/jwt_secret/credential_value/private_key are stripped.",
            "Failed-login gotcha: HTTPException rolls back the request session → record_login_event() uses its OWN session that commits immediately.",
            "X-Request-Id response header correlates every mutation to its originating request end-to-end.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_observability_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Observability",
        "Centralized logs · system health · built-in analytics",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Centralized logs: collectors ship filtered container logs to central → ClickHouse logs table.",
            "logs table uses toDateTime() wrapper in TTL (required for CH 24.3 DateTime64).",
            "log_level_filter per Collector row controls minimum shipped level — operator-tunable.",
            "System Health page: per-subsystem status (Postgres/CH/Redis/scheduler/collectors/etc.) excluding alert state.",
            "Top-bar health indicator dot hits a cached /system/health/status (admin-only).",
            "Metabase on central1 (exposed via HAProxy /metabase) — ad-hoc SQL over ClickHouse result tables.",
            "Grafana on central4 — operator dashboards over the same data.",
            "Built-in Analytics UI: SQL Explorer + custom dashboards with react-grid-layout + cross-widget variables.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_deployment_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Deployment & operations",
        "Docker Compose per role · one-command ./deploy.sh",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Central: central1-3 run \"central-ha\" compose (app + Patroni + Redis + sentinel + Metabase on c1). central4 runs \"central\" + Grafana.",
            "ClickHouse on central3+4 · Keeper quorum voter on central1 · etcd on central1-3.",
            "Workers: single compose project per collector node (/opt/monctl/collector/).",
            "./deploy.sh central builds once, saves the image, scp's + docker load on all 4 nodes in parallel (~20s).",
            "./deploy.sh collector 31 32 — deploy to a subset of hosts when you need to roll forward carefully.",
            "Alembic migrations run on central startup (lifespan hook); create_all only runs on fresh DBs (PR #20 fix).",
            "Everything observable via standard docker compose logs + the built-in log collection.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


def make_status_slide(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)
    _add_header(
        slide, "Status & roadmap",
        "What's live today and what's next",
    )
    _add_bullets(
        slide, Inches(0.7), Inches(1.9), Inches(12), Inches(5),
        [
            "Live in prod: full central HA · 4 workers · SNMP + ping + HTTP + SSH app fleet · replicated ClickHouse.",
            "Event policy rework COMPLETE (PRs #20–#33). Legacy EventEngine / EventPolicy / ActiveEvent retired.",
            "Incident engine features shipped: severity ladders · scope filters · dependencies · flap guard · companion-clear (any/all).",
            "Dashboard editor · analytics dashboards · custom SQL explorer · automation run books.",
            "Discovery auto-rename · bulk import · hierarchical templates · device-type resolvers.",
            "Deferred: phase 2c (string-template correlation groups) — no operator demand yet.",
            "In-flight: connector alias → slot model (branch pushed, pending migration rebase).",
            "Future: additional SNMP apps for new device families as they enter the fleet.",
        ],
        size=15,
    )
    _add_footer(slide, idx, total)


# ── Main ───────────────────────────────────────────────────
def build(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        make_title_slide,
        make_problem_slide,
        make_architecture_slide,
        make_stack_slide,
        make_central_slide,
        make_collector_slide,
        make_ha_slide,
        make_storage_slide,
        make_scheduling_slide,
        make_apps_slide,
        make_credentials_slide,
        make_interface_slide,
        make_alerting_slide,
        make_incidents_slide,
        make_automations_slide,
        make_security_slide,
        make_audit_slide,
        make_observability_slide,
        make_deployment_slide,
        make_status_slide,
    ]
    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        builder(prs, i, total)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"wrote {output} ({total} slides)")


if __name__ == "__main__":
    build(Path(__file__).resolve().parent.parent / "MonCTL-technical-overview.pptx")
